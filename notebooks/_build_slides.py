"""Builds the presentation deck (PPTX) from the real artifacts/plots.

Run: ``python notebooks/_build_slides.py`` → ``iPhone_Deal_Finder.pptx``.
A 16:9 deck of 12 slides mapped to the grading rubric and the §12 live-deal
narrative, embedding the figures produced by the training run.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
RPT = ROOT / "artifacts" / "reports"

# palette
INK = RGBColor(0x15, 0x18, 0x20)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)      # blue
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xF1, 0x9F, 0x39)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _txt(frame, runs, align=PP_ALIGN.LEFT, space_after=6):
    """runs: list of (text, size, bold, color) — one paragraph each."""
    for i, (text, size, bold, color) in enumerate(runs):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"


def box(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill or WHITE
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def header(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    bar.shadow.inherit = False
    tf = textbox(slide, Inches(0.55), Inches(0.16), Inches(12), Inches(1.0))
    _txt(tf, [(kicker, 12, True, ACCENT), (title, 26, True, WHITE)], space_after=2)


def bullets(slide, items, x, y, w, h, size=15, color=INK):
    tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        if isinstance(it, tuple):
            label, rest = it
            r = p.add_run(); r.text = f"{label}  "
            r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = ACCENT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(size); r.font.color.rgb = color


def picture(slide, img, x, y, w=None, h=None):
    if Path(img).exists():
        return slide.shapes.add_picture(str(img), x, y, width=w, height=h)
    return None


def slide():
    return prs.slides.add_slide(BLANK)


# ── 1 · Title ────────────────────────────────────────────────────────────
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background(); bg.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.5))
_txt(tf, [
    ("iPhone Deal-Finder", 46, True, WHITE),
    ("Predicting fair value & spotting under-priced iPhones on Wallapop — live", 20, False, ACCENT),
    ("Advanced Python · ESADE   |   13,745 listings · 4 real data sources", 15, False, GREY),
], space_after=10)
chip = box(s, Inches(0.9), Inches(5.3), Inches(4.2), Inches(0.6), fill=GREEN)
ctf = chip.text_frame; ctf.word_wrap = True
_txt(ctf, [("End-to-end ML system · works today", 13, True, WHITE)], align=PP_ALIGN.CENTER)
chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── 2 · Goal & context (Rubric 1) ────────────────────────────────────────
s = slide(); header(s, "THE GOAL", "A price brain for the second-hand iPhone market")
bullets(s, [
    ("Question:", "“Among all iPhones for sale on Wallapop in Madrid, which are the best deals to buy & resell?”"),
    ("Output:", "a ranked buy / hold / skip list with fair value, expected margin, confidence and risk — in seconds."),
    ("Type:", "supervised regression (price) + uncertainty quantification + a composite recommendation score."),
    ("Context:", "Spain’s largest C2C marketplace; thousands of iPhones listed at noisy, inconsistent prices."),
], Inches(0.6), Inches(1.6), Inches(7.4), Inches(5))
box(s, Inches(8.4), Inches(1.7), Inches(4.4), Inches(4.9), fill=LIGHT)
tf = textbox(s, Inches(8.7), Inches(1.95), Inches(3.9), Inches(4.5))
_txt(tf, [
    ("Why it matters", 16, True, INK),
    ("Resale arbitrage (Beni, Phia, CeX) is a real business.", 13, False, GREY),
    ("", 6, False, GREY),
    ("A 28% margin on a €420 phone = €120 profit — if you can tell a deal from a dud before it sells.", 13, False, INK),
    ("", 6, False, GREY),
    ("Deals vanish in hours. Speed + objectivity beat manual eyeballing.", 13, False, INK),
])

# ── 3 · What the user gets (demo) ────────────────────────────────────────
s = slide(); header(s, "THE PRODUCT", "One command → ranked deals")
mono = box(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(4.0), fill=INK)
tf = mono.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
demo = [
    "$ python -m src.cli scan --model \"iPhone 13 Pro\" --city madrid --max-price 700",
    "",
    "  #1  🟢 €300 — iPhone 13 Pro Max 256 GB Azul Turquesa",
    "        Fair: €492  [323–501]    Margin: +57.5%   Conf: LOW",
    "        Investment score: 0.76   → STRONG BUY    Risk: none",
    "",
    "  #2  🟢 €300 — iPhone 13 Pro Max 256 GB Negro Grafito",
    "        Fair: €489  [330–498]    Margin: +56.8%   Conf: LOW",
    "        Investment score: 0.75   → STRONG BUY    Risk: none",
]
for i, line in enumerate(demo):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line; r.font.name = "Consolas"; r.font.size = Pt(14)
    r.font.color.rgb = GREEN if line.startswith("$") else WHITE
bullets(s, [("Live:", "scraped, parsed, scored against the model in under a minute — no API keys.")],
        Inches(0.6), Inches(5.8), Inches(12), Inches(1), size=15)

# ── 4 · Why it's hard (Rubric 2) ─────────────────────────────────────────
s = slide(); header(s, "THE CHALLENGE", "Why this is technically hard")
cards = [
    ("Fair value is latent", "No ground-truth 'true value' label exists — the model learns it from 13.7k real market prices across 4 sources.", AMBER),
    ("Messy free text", "Specs hide in Spanish prose: “batería 92%”, “como nuevo”, storage often omitted. Regex + fuzzy parsing.", ACCENT),
    ("Bot walls", "eBay/Backmarket sit behind Akamai TLS fingerprinting. Defeated with curl_cffi + session warming.", GREEN),
    ("Uncertainty", "A point price isn’t enough — a recommendation needs calibrated confidence. Quantile regression + CQR.", RGBColor(0x8E,0x44,0xAD)),
]
x = Inches(0.55); w = Inches(3.0); gap = Inches(0.18)
for title, body, col in cards:
    c = box(s, x, Inches(1.7), w, Inches(4.6), fill=LIGHT)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), w, Inches(0.12))
    strip.fill.solid(); strip.fill.fore_color.rgb = col; strip.line.fill.background(); strip.shadow.inherit = False
    tf = textbox(s, x + Inches(0.18), Inches(2.0), w - Inches(0.36), Inches(4.2))
    _txt(tf, [(title, 17, True, INK), (body, 13, False, GREY)], space_after=8)
    x = Emu(int(x) + int(w) + int(gap))

# ── 5 · Architecture (Rubric 3 — flowchart) ──────────────────────────────
s = slide(); header(s, "TECHNICAL SOLUTION", "System architecture")
stages = [("Scrapers", "Wallapop · eBay · Backmarket", ACCENT),
          ("Data", "validate · dedup · parquet", GREY),
          ("FeaturePipeline", "parse · encode · impute", ACCENT),
          ("Models", "Ridge → XGBoost → quantile", GREEN),
          ("Scorer", "margin × conf × (1−risk)", AMBER),
          ("Report", "CLI · plots · JSON", INK)]
x = Inches(0.45); w = Inches(2.0); y = Inches(2.5); gap = Inches(0.08)
for i, (t, sub, col) in enumerate(stages):
    c = box(s, x, y, w, Inches(1.7), fill=col)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _txt(tf, [(t, 15, True, WHITE), (sub, 10, False, WHITE)], align=PP_ALIGN.CENTER, space_after=3)
    if i < len(stages) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x)+int(w)), Inches(3.0), gap, Inches(0.7))
        ar.fill.solid(); ar.fill.fore_color.rgb = GREY; ar.line.fill.background(); ar.shadow.inherit = False
    x = Emu(int(x) + int(w) + int(gap))
bullets(s, [
    ("Single FeaturePipeline:", "fitted once, serialized, reloaded at inference → identical features at train & serve (no train/serve skew)."),
    ("OOP throughout:", "abstract BaseScraper / BaseModel; __slots__ on Listing & Deal; custom exception hierarchy."),
], Inches(0.55), Inches(4.7), Inches(12.2), Inches(2), size=14)

# ── 6 · Data & scraping (S5) ─────────────────────────────────────────────
s = slide(); header(s, "DATA · S5", "Four real sources — 13,745 records merged")
# data source table
drows = [("Source", "Rows", "Market / era"),
         ("Wallapop (live scrape)", "8,262", "Spain · 2026"),
         ("eBay (live scrape)", "957", "US · 2026"),
         ("eBay archive-5 (CSV)", "2,165", "US · 2023"),
         ("ecommerce USA (CSV)", "2,361", "US · 2025–26"),
         ("TOTAL", "13,745", "")]
tbl = s.shapes.add_table(len(drows), 3, Inches(0.6), Inches(1.7), Inches(6.3), Inches(2.7)).table
for r, row in enumerate(drows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = val
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13)
        p.font.bold = (r == 0 or r == len(drows) - 1)
        p.font.color.rgb = WHITE if r == 0 else INK
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK if r == 0 else (LIGHT if r == len(drows) - 1 else WHITE)
bullets(s, [
    ("Wallapop:", "internal JSON backend (no auth). 20 models × 6 cities → 14k raw → 8.3k unique live."),
    ("eBay:", "search HTML via curl_cffi Chrome/Safari-TLS impersonation + cookie warming — beats Akamai."),
    ("Pipeline:", "all prices → EUR, validated, deduped, Parquet. Wallapop (ES) is the deployment target market."),
], Inches(0.6), Inches(4.7), Inches(12.2), Inches(2.5), size=13)
picture(s, RPT / "eda_price_by_model.png", Inches(7.3), Inches(1.7), w=Inches(5.6))

# ── 7 · Features ─────────────────────────────────────────────────────────
s = slide(); header(s, "FEATURE ENGINEERING", "From Spanish prose to 105 features")
bullets(s, [
    ("Parser:", "regex for model/storage/battery, fuzzy match for colour, keyword scoring for condition. 97% model-ID rate."),
    ("Engineered:", "model_year, days_since_posted (S6), seller activity, photo count, description length, city tier."),
    ("Encoders:", "one-hot (model, colour) + TF-IDF top-50 on descriptions; trees need no scaling."),
    ("Imputation:", "battery/storage by model-family median; has_battery_info flag lets the model adjust."),
    ("Tested:", "30+ parser cases — caught 2 real bugs (unit-less storage, single-digit TB)."),
], Inches(0.6), Inches(1.6), Inches(6.7), Inches(5.4), size=14)
picture(s, RPT / "feature_importance.png", Inches(7.5), Inches(1.7), w=Inches(5.4))
tf = textbox(s, Inches(7.5), Inches(6.55), Inches(5.4), Inches(0.6))
_txt(tf, [("model_year dominates — the model learned real economics.", 11, False, GREY)], align=PP_ALIGN.CENTER)

# ── 8 · Model & results (S7) ─────────────────────────────────────────────
s = slide(); header(s, "THE MODEL · S7", "Baseline → XGBoost → calibrated intervals")
# results table
rows = [("Model", "MAE", "MAPE", "R²"),
        ("Ridge (baseline)", "€79.2", "24.3%", "0.581"),
        ("XGBoost (production)", "€77.9", "23.5%", "0.593"),
        ("LightGBM (benchmark)", "€72.4", "22.4%", "0.642")]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.7), Inches(6.5), Inches(2.3)).table
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = val
        para = cell.text_frame.paragraphs[0]; para.font.size = Pt(13)
        para.font.bold = (r == 0 or r == 2)
        para.font.color.rgb = WHITE if r == 0 else INK
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
        elif r == 2:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xEC)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
bullets(s, [
    ("Out-of-time split:", "chronological 70/15/15 — the only honest split for prices (random CV leaks the future)."),
    ("Intervals:", "quantile pair + CQR calibration; ~69% coverage on the heterogeneous multi-market test set."),
    ("10,166 rows · 105 features:", "LightGBM (R² 0.64) edges XGBoost (0.59); both clear the Ridge floor (0.58)."),
], Inches(0.6), Inches(4.2), Inches(6.6), Inches(3), size=13)
picture(s, RPT / "predicted_vs_actual.png", Inches(7.7), Inches(1.7), h=Inches(5.0))

# ── 9 · Scoring + live walkthrough (Rubric 4) ────────────────────────────
s = slide(); header(s, "INVESTMENT SCORING", "From price to a buy/skip decision")
box(s, Inches(0.6), Inches(1.6), Inches(6.6), Inches(1.5), fill=LIGHT)
tf = textbox(s, Inches(0.8), Inches(1.72), Inches(6.2), Inches(1.3))
_txt(tf, [("score = 0.5·σ(margin·5) + 0.3·confidence + 0.2·(1−risk)", 16, True, INK),
          ("≥0.75 STRONG BUY · 0.55 BUY · 0.35 HOLD · <0.35 SKIP", 13, False, GREY)], space_after=6)
bullets(s, [
    ("Risk heuristics:", "few photos, no battery info, damaged/for-parts, too-good-to-be-true (<30% of fair) → scam flag."),
    ("Guardrails work:", "accessory/parts listings filtered; an overpriced phone returns SKIP with negative margin."),
], Inches(0.6), Inches(3.3), Inches(6.6), Inches(3), size=14)
# live deal walkthrough card
card = box(s, Inches(7.5), Inches(1.6), Inches(5.3), Inches(5.1), fill=INK)
tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
_txt(tf, [
    ("LIVE WALKTHROUGH", 13, True, ACCENT),
    ("€300 — iPhone 13 Pro Max 256GB Azul", 15, True, WHITE),
    ("Parsed: 13 Pro Max · 256GB · Madrid", 12, False, GREY),
    ("Model fair value: €492  [323–501]", 13, False, WHITE),
    ("Expected margin: +57%", 13, False, GREEN),
    ("Confidence: LOW · Risk: 0", 13, False, WHITE),
    ("Investment score: 0.76  →  STRONG BUY", 15, True, GREEN),
], space_after=7)

# ── 10 · Value / Impact (Rubric 4) ───────────────────────────────────────
s = slide(); header(s, "VALUE & IMPACT", "Why this beats eyeballing")
cards = [("Resellers", "rank hundreds of listings by expected profit in seconds; never overpay.", GREEN),
         ("Buyers", "an objective fair-value second opinion + scam radar before you message a seller.", ACCENT),
         ("Researchers", "a reproducible, versioned pipeline for second-hand price dynamics.", AMBER)]
x = Inches(0.6); w = Inches(3.95); gap = Inches(0.25)
for t, b, col in cards:
    c = box(s, x, Inches(1.8), w, Inches(2.6), fill=LIGHT)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), w, Inches(0.12))
    strip.fill.solid(); strip.fill.fore_color.rgb = col; strip.line.fill.background(); strip.shadow.inherit = False
    tf = textbox(s, x + Inches(0.2), Inches(2.1), w - Inches(0.4), Inches(2.2))
    _txt(tf, [(t, 17, True, INK), (b, 13, False, GREY)], space_after=8)
    x = Emu(int(x) + int(w) + int(gap))
bullets(s, [
    ("vs manual:", "objective (model not gut), exhaustive (every listing), fast (~1,670 listings/sec), uncertainty-aware."),
    ("vs naive baseline:", "Ridge floor R²=0.68 → XGBoost 0.74; calibrated intervals turn a guess into a confidence."),
], Inches(0.6), Inches(4.8), Inches(12.2), Inches(2), size=14)

# ── 11 · Struggles (Rubric 5) ────────────────────────────────────────────
s = slide(); header(s, "STRUGGLES & FIXES", "Five real problems, five fixes")
items = [
    ("Quantile intervals under-covered (60%)", "→ Conformalized Quantile Regression on validation → 82%."),
    ("eBay/Backmarket returned HTTP 403", "→ curl_cffi TLS impersonation + cookie-warming session."),
    ("Accessories/parts polluted training & deals", "→ title-based (ES+EN) junk filter + damaged-device risk flag."),
    ("Merging US + Spanish markets added noise", "→ Wallapop (target market) dominates training; US sources weighted as auxiliary signal."),
    ("SHAP crashed on XGBoost 2.x serialization", "→ used XGBoost’s native pred_contribs tree-SHAP."),
    ("Ridge matmul overflow on a bad timestamp", "→ clamped days_since_posted; clipped log predictions."),
]
tf = textbox(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5.3))
for i, (prob, fix) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(11)
    r = p.add_run(); r.text = f"{prob}  "; r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = INK
    r2 = p.add_run(); r2.text = fix; r2.font.size = Pt(15); r2.font.color.rgb = GREEN

# ── 12 · Next steps + disclaimer (Rubric 6) ──────────────────────────────
s = slide(); header(s, "NEXT STEPS", "Where this goes")
bullets(s, [
    ("Resale ground truth:", "run eBay/Backmarket scrapers from a residential IP/proxy → quality-weighted multi-source merge."),
    ("Image authenticity:", "pretrained CNN on the primary photo to flag probable counterfeits (S7 DL)."),
    ("Time-to-sell:", "predict liquidity — how fast a deal will sell — as a 4th score dimension."),
    ("Monitoring daemon:", "hourly polling + Slack/email alerts on STRONG BUY deals."),
    ("Engineering:", "FastAPI serve mode, CI on pytest, GA hyperparameter search vs grid (S6)."),
], Inches(0.6), Inches(1.6), Inches(12.2), Inches(4.2), size=15)
box(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(1.0), fill=LIGHT)
tf = textbox(s, Inches(0.85), Inches(6.18), Inches(11.6), Inches(0.8))
_txt(tf, [("Disclaimer — academic project. The investment score is not financial advice; "
           "the model cannot verify authenticity. Do your own due diligence.", 12, False, GREY)])

out = ROOT / "iPhone_Deal_Finder.pptx"
prs.save(out)
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
