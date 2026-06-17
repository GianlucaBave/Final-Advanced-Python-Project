"""Build the Apple/glassmorphism .pptx deck (SF Pro, dark gradients, glass cards).

Run: python notebooks/_build_apple_deck.py  ->  iPhone_Deal_Finder.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
RPT = ROOT / "artifacts" / "reports"
ASSETS = ROOT / "assets"

# palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0B, 0x0B, 0x0F)
SUB = RGBColor(0xAE, 0xAE, 0xB2)        # secondary label
BLUE = RGBColor(0x0A, 0x84, 0xFF)
GREEN = RGBColor(0x30, 0xD1, 0x58)
AMBER = RGBColor(0xFF, 0x9F, 0x0A)
RED = RGBColor(0xFF, 0x45, 0x3A)
PURPLE = RGBColor(0xBF, 0x5A, 0xFF)
HEAD = "SF Pro Display"
BODY = "SF Pro Text"
MONO = "SF Mono"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- low-level helpers ----------
def _alpha(shape, rgb: RGBColor, a_pct: int):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb
    srgb = shape._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    el = srgb.makeelement(qn("a:alpha"), {"val": str(int(a_pct * 1000))}); srgb.append(el)


def _gradient(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(t)):
            spPr.remove(el)
    g = spPr.makeelement(qn("a:gradFill"), {}); lst = g.makeelement(qn("a:gsLst"), {})
    for pos, c in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        sc = gs.makeelement(qn("a:srgbClr"), {"val": c}); gs.append(sc); lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    ln.addprevious(g) if ln is not None else spPr.append(g)


def _no_line(shape): shape.line.fill.background()
def _noshadow(shape): shape.shadow.inherit = False


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _no_line(bg); _noshadow(bg)
    _gradient(bg, [(0, "12121A"), (55, "0B0B0F"), (100, "16121F")], ang=115)
    return s


def glow(s, cx, cy, d, hexc, a=14):
    e = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), d, d)
    _no_line(e); _noshadow(e); _alpha(e, RGBColor.from_string(hexc), a)
    return e


def glass(s, x, y, w, h, a=9, radius=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    _alpha(shp, WHITE, a)
    shp.line.color.rgb = RGBColor(0x9A, 0x9A, 0xB0); shp.line.width = Pt(1.0)
    _noshadow(shp)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, r in enumerate(runs):
        txt, size, bold, col = r[0], r[1], r[2], r[3]
        font = r[4] if len(r) > 4 else BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        run = p.add_run(); run.text = txt
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = col; run.font.name = font
    return tb


def bullets(s, x, y, w, h, items, size=15, col=WHITE, lead=None, gap=9):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            r1 = p.add_run(); r1.text = it[0] + "  "
            r1.font.bold = True; r1.font.size = Pt(size); r1.font.color.rgb = lead or BLUE; r1.font.name = BODY
            r2 = p.add_run(); r2.text = it[1]
            r2.font.size = Pt(size); r2.font.color.rgb = col; r2.font.name = BODY
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(size); r.font.color.rgb = col; r.font.name = BODY
    return tb


def pill(s, x, y, w, label, color, txtcol=WHITE):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    _alpha(p, color, 90); _no_line(p); _noshadow(p)
    try: p.adjustments[0] = 0.5
    except Exception: pass
    tf = p.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = label; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = txtcol; r.font.name = BODY
    return p


def kicker(s, txt, color=BLUE):
    text(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4),
         [(txt.upper(), 13, True, color, BODY)], sp=0)


def title(s, txt, y=0.95, size=34):
    text(s, Inches(0.7), Inches(y), Inches(12), Inches(1.0), [(txt, size, True, WHITE, HEAD)], sp=0)


def pic(s, name, x, y, w=None, h=None):
    p = RPT / name
    if p.exists():
        return s.shapes.add_picture(str(p), x, y, width=w, height=h)


def logo(s, name, x, y, w=None, h=None):
    p = ASSETS / name
    if p.exists():
        return s.shapes.add_picture(str(p), x, y, width=w, height=h)


def iphone(s, x, y, h, screen_hex="0A84FF"):
    """A minimal glass iPhone silhouette."""
    w = Emu(int(h * 0.49))
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _alpha(body, WHITE, 10); body.line.color.rgb = RGBColor(0xAA, 0xAA, 0xC0); body.line.width = Pt(1.2)
    _noshadow(body)
    try: body.adjustments[0] = 0.16
    except Exception: pass
    inset = Emu(int(h * 0.025))
    scr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x) + int(inset)), Emu(int(y) + int(inset)),
                             Emu(int(w) - 2 * int(inset)), Emu(int(h) - 2 * int(inset)))
    _no_line(scr); _noshadow(scr)
    _gradient(scr, [(0, screen_hex), (100, "1C1C28")], ang=120)
    try: scr.adjustments[0] = 0.14
    except Exception: pass
    # notch
    nw = Emu(int(w * 0.34)); nh = Emu(int(h * 0.03))
    notch = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x) + (int(w) - int(nw)) // 2),
                               Emu(int(y) + int(inset) + int(h * 0.012)), nw, nh)
    _alpha(notch, INK, 85); _no_line(notch); _noshadow(notch)
    return body


# ════════════════════════ SLIDES ════════════════════════
# 1 · TITLE
s = slide()
glow(s, int(SW * 0.78), int(SH * 0.22), Inches(7), "0A84FF", 16)
glow(s, int(SW * 0.15), int(SH * 0.85), Inches(6), "BF5AF2", 12)
iphone(s, Inches(9.7), Inches(1.5), Inches(4.4), "0A84FF")
logo(s, "apple_white.png", Inches(10.55), Inches(2.05), w=Inches(0.46))  # Apple mark on screen
# deal card overlay on phone
dc = glass(s, Inches(8.55), Inches(3.0), Inches(3.4), Inches(1.5), a=14)
text(s, Inches(8.75), Inches(3.12), Inches(3.0), Inches(1.3), [
    ("€299  iPhone 13 Pro", 14, True, WHITE, BODY),
    ("Fair value €409", 12, False, SUB, BODY),
    ("+31% margin · score 0.76", 12, True, GREEN, BODY),
], sp=4)
pill(s, Inches(8.75), Inches(4.05), Inches(1.7), "STRONG BUY", GREEN, INK)
text(s, Inches(0.8), Inches(2.2), Inches(8.3), Inches(3), [
    ("iPhone Deal-Finder", 54, True, WHITE, HEAD),
    ("Machine learning that predicts an iPhone's fair value and", 20, False, SUB, BODY),
    ("spots underpriced listings, on live data, in seconds.", 20, False, SUB, BODY),
], sp=8)
pill(s, Inches(0.85), Inches(5.45), Inches(4.6), "Advanced Python · ESADE · Final Project", BLUE)
text(s, Inches(0.85), Inches(6.2), Inches(8), Inches(0.6),
     [("Regression  ·  Web scraping  ·  XGBoost  ·  13,745 real listings", 14, False, SUB, MONO)], sp=0)

# 2 · GOAL (Rubric 1)
s = slide(); glow(s, int(SW*0.85), int(SH*0.2), Inches(6), "0A84FF", 12)
kicker(s, "01 · Project Explanation /10"); title(s, "The goal: a price brain for second-hand iPhones")
g = glass(s, Inches(0.7), Inches(1.7), Inches(7.3), Inches(5.1))
bullets(s, Inches(1.0), Inches(2.0), Inches(6.8), Inches(4.6), [
    ("Question", "“Of all iPhones for sale right now, which are the best deals to buy and resell?”"),
    ("Model", "supervised regression that predicts fair market value, plus calibrated uncertainty."),
    ("On top", "an investment score → a clear BUY / HOLD / SKIP decision with reasons."),
    ("Context", "Wallapop is Spain's largest second-hand marketplace; prices are noisy and inconsistent."),
    ("Real use", "resale arbitrage and smart buying: rank live listings by expected profit, flag scams."),
], size=16, gap=14)
g2 = glass(s, Inches(8.3), Inches(1.7), Inches(4.3), Inches(5.1), a=7)
text(s, Inches(8.6), Inches(2.0), Inches(3.7), Inches(1), [("Why it matters", 18, True, WHITE, HEAD)], sp=4)
text(s, Inches(8.6), Inches(2.7), Inches(3.8), Inches(4), [
    ("A 28% margin on a €420 phone is €120 profit.", 14, False, WHITE, BODY),
    ("if you can tell a deal from a dud before it sells.", 14, False, SUB, BODY),
    ("", 8, False, SUB, BODY),
    ("Deals vanish in hours.", 14, True, WHITE, BODY),
    ("Speed + objectivity beat manual eyeballing.", 14, False, SUB, BODY),
], sp=6)

# 3 · PRODUCT DEMO
s = slide(); glow(s, int(SW*0.2), int(SH*0.3), Inches(6), "30D158", 10)
kicker(s, "The product", GREEN); title(s, "One command → ranked deals")
card = glass(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(4.9), a=8)
text(s, Inches(1.0), Inches(1.95), Inches(11), Inches(0.5),
     [("$ python -m src.cli scan --model \"iPhone 13 Pro\" --city madrid --max-price 700", 15, True, GREEN, MONO)], sp=0)
# two deal cards
for i, (x, price, model, fair, margin, conf, score, dec, col) in enumerate([
    (0.95, "€299", "iPhone 13 Pro · TIENDA/GARANTÍA", "€409  [300–403]", "+31.0%", "MEDIUM", "0.76", "STRONG BUY", GREEN),
    (6.85, "€300", "iPhone 13 Pro Max 256GB Titanio", "€569  [347–596]", "+82.2%", "LOW", "0.73", "BUY", AMBER)]):
    c = glass(s, Inches(x), Inches(2.7), Inches(5.6), Inches(3.5), a=12)
    text(s, Inches(x+0.3), Inches(2.9), Inches(5.0), Inches(2.9), [
        (f"{price}", 30, True, WHITE, HEAD),
        (model, 13, False, SUB, BODY),
        ("", 6, False, SUB, BODY),
        (f"Fair value: {fair}", 15, True, WHITE, BODY),
        (f"Margin {margin}   ·   Confidence {conf}", 14, False, SUB, BODY),
        (f"Investment score {score}", 15, True, (GREEN if dec=="STRONG BUY" else AMBER), BODY),
    ], sp=5)
    pill(s, Inches(x+0.3), Inches(5.65), Inches(2.0), dec, col, INK)

# 4 · CHALLENGE (Rubric 2)
s = slide(); kicker(s, "02 · Challenge Definition /15"); title(s, "Why this is genuinely hard")
cards = [
    ("No labels", "Wallapop shows asking prices, never sold prices, so the model learns fair value from the crowd.", AMBER),
    ("Messy free text", "Specs hide in Spanish prose: “batería 92%”, “como nuevo”, storage often omitted.", BLUE),
    ("Bot walls", "eBay & Backmarket fingerprint the TLS handshake and 403 datacenter IPs.", GREEN),
    ("Uncertainty", "A single price isn't enough; a recommendation needs calibrated confidence.", PURPLE),
]
x = Inches(0.7); w = Inches(2.95); gap = Inches(0.12)
for t, b, c in cards:
    g = glass(s, x, Inches(1.9), w, Inches(4.6), a=8)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.25), Inches(2.2), Inches(0.5), Inches(0.12))
    _alpha(strip, c, 100); _no_line(strip); _noshadow(strip)
    text(s, x+Inches(0.25), Inches(2.45), w-Inches(0.5), Inches(4), [(t, 18, True, WHITE, HEAD), ("", 6, False, SUB, BODY), (b, 14, False, SUB, BODY)], sp=8)
    x = Emu(int(x)+int(w)+int(gap))

# 5 · DATA
s = slide(); glow(s, int(SW*0.8), int(SH*0.8), Inches(6), "BF5AF2", 10)
kicker(s, "Data strategy"); title(s, "13,745 real listings · 4 marketplaces")
logo(s, "ebay.png", Inches(11.3), Inches(0.5), w=Inches(1.3))      # marketplace marks
logo(s, "apple_white.png", Inches(10.6), Inches(0.5), w=Inches(0.4))
pic(s, "slide_data_sources.png", Inches(0.6), Inches(1.7), h=Inches(5.3))
g = glass(s, Inches(7.0), Inches(1.8), Inches(5.7), Inches(4.9), a=8)
bullets(s, Inches(7.3), Inches(2.05), Inches(5.2), Inches(4.5), [
    ("Wallapop · 8,262", "scraped LIVE from its JSON backend (no API, no keys). 20 models × 6 Spanish cities. Deployment market."),
    ("eBay sold · 957", "real completed-sale prices with dates, via sold-listings HTML."),
    ("eBay archive · 2,165", "historical 2023 US listings for temporal depth."),
    ("US e-commerce · 2,361", "2026 resale records, sale-evidenced."),
    ("Normalized", "all prices → EUR, stored in Parquet, quality-weighted at merge."),
], size=14.5, gap=12)

# 6 · ARCHITECTURE (Rubric 3)
s = slide(); kicker(s, "03 · Technical Solution /25", BLUE); title(s, "System architecture")
stages = [("Scrapers", "Wallapop · eBay", BLUE), ("Data", "validate · parquet", SUB),
          ("FeaturePipeline", "parse · encode", BLUE), ("Models", "XGBoost · quantile", GREEN),
          ("Scorer", "margin × conf", AMBER), ("Report", "CLI · plots", PURPLE)]
x = Inches(0.55); w = Inches(1.95); y = Inches(2.4)
for i, (t, sub, c) in enumerate(stages):
    g = glass(s, x, y, w, Inches(1.7), a=10)
    text(s, x, y+Inches(0.45), w, Inches(1), [(t, 14.5, True, WHITE, HEAD), (sub, 10.5, False, SUB, BODY)], align=PP_ALIGN.CENTER, sp=3)
    if i < len(stages)-1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Emu(int(x)+int(w)-Inches(0.02)), y+Inches(0.62), Inches(0.22), Inches(0.45))
        _alpha(ar, RGBColor.from_string("0A84FF"), 80); _no_line(ar); _noshadow(ar)
    x = Emu(int(x)+int(w)+Inches(0.14))
g = glass(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.1), a=7)
bullets(s, Inches(1.0), Inches(4.85), Inches(11.3), Inches(1.7), [
    ("Single FeaturePipeline", "fitted once, serialized, reloaded at inference → identical features at train & serve (zero skew)."),
    ("Clean OOP", "abstract BaseScraper / BaseModel, __slots__ containers, custom exception hierarchy, typed config, structured logging."),
], size=14.5, gap=10)

# 7 · FEATURES
s = slide(); kicker(s, "Feature engineering"); title(s, "From Spanish prose to 105 features")
g = glass(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(4.9), a=8)
text(s, Inches(1.0), Inches(2.05), Inches(5.0), Inches(0.5), [("raw listing", 13, True, SUB, MONO)], sp=2)
text(s, Inches(1.0), Inches(2.45), Inches(5.0), Inches(0.8), [("“iPhone 13 Pro 256GB Sierra Blue batería 92% como nuevo”", 14, True, WHITE, BODY)], sp=2)
text(s, Inches(1.0), Inches(3.5), Inches(5.0), Inches(0.4), [("↓  parsed", 13, True, GREEN, MONO)], sp=2)
bullets(s, Inches(1.0), Inches(3.95), Inches(5.0), Inches(2.6), [
    ("model", "iPhone 13 Pro · year 2021"), ("storage", "256 GB"),
    ("battery", "92%"), ("colour", "sierra blue"), ("condition", "like-new (4/5)"),
], size=14, lead=BLUE, gap=7)
g2 = glass(s, Inches(6.6), Inches(1.8), Inches(6.0), Inches(4.9), a=8)
bullets(s, Inches(6.9), Inches(2.1), Inches(5.5), Inches(4.5), [
    ("Regex + fuzzy", "model, storage, battery, colour, condition from free text (97% model-ID rate)."),
    ("TF-IDF", "top-50 description tokens capture wording signal."),
    ("Engineered", "days-since-posted, seller activity, photo count, city tier."),
    ("Market / source", "indicator features let one model serve multiple marketplaces."),
    ("Imputation", "battery/storage by model-family median; nothing dropped silently."),
], size=14.5, gap=12)

# 8 · MODEL (Rubric 3 cont.)
s = slide(); kicker(s, "03 · Technical Solution /25", BLUE); title(s, "Baseline → boosting → calibrated intervals")
pic(s, "slide_model_comparison.png", Inches(0.6), Inches(1.85), w=Inches(6.0))
pic(s, "slide_pred_vs_actual.png", Inches(6.9), Inches(1.7), h=Inches(4.4))
g = glass(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(1.05), a=8)
text(s, Inches(1.0), Inches(6.18), Inches(11.3), Inches(0.8), [
    ("Out-of-time split (no future leakage)  ·  LightGBM best: R² 0.642, MAE €72  ·  80% prediction intervals via Conformalized Quantile Regression", 14, True, WHITE, BODY)], sp=0)

# 9 · INVESTMENT SCORE
s = slide(); glow(s, int(SW*0.8), int(SH*0.25), Inches(6), "30D158", 10)
kicker(s, "The recommendation engine", GREEN); title(s, "Investment score → decision")
g = glass(s, Inches(0.7), Inches(1.8), Inches(7.2), Inches(2.0), a=10)
text(s, Inches(1.0), Inches(2.05), Inches(6.7), Inches(1.6), [
    ("investment_score =", 16, True, WHITE, MONO),
    ("0.5 · σ(margin·5)  +  0.3 · confidence  +  0.2 · (1 − risk)", 17, True, BLUE, MONO),
], sp=8)
for i, (lbl, col, rng) in enumerate([("STRONG BUY", GREEN, "≥ 0.75"), ("BUY", RGBColor(0x9B,0xE8,0xA8), "0.55–0.75"),
                                      ("HOLD", SUB, "0.35–0.55"), ("SKIP", RED, "< 0.35")]):
    y = Inches(4.05 + i*0.62)
    pill(s, Inches(0.85), y, Inches(2.0), lbl, col, INK if lbl in ("STRONG BUY", "BUY", "HOLD") else WHITE)
    text(s, Inches(3.0), y, Inches(2), Inches(0.42), [(rng, 14, True, WHITE, MONO)], anchor=MSO_ANCHOR.MIDDLE, sp=0)
g2 = glass(s, Inches(8.2), Inches(1.8), Inches(4.4), Inches(4.9), a=8)
text(s, Inches(8.5), Inches(2.05), Inches(3.8), Inches(0.5), [("Risk heuristics", 18, True, WHITE, HEAD)], sp=4)
bullets(s, Inches(8.5), Inches(2.7), Inches(3.9), Inches(3.9), [
    "Price < 30% of fair → scam flag", "No battery info", "Fewer than 2 photos",
    "“para piezas” / damaged wording", "iCloud-locked / suspicious",
], size=14, gap=11)

# 10 · VALUE (Rubric 4)
s = slide(); kicker(s, "04 · Value or Impact /15", BLUE); title(s, "Why it beats eyeballing")
for i, (t, b, c) in enumerate([
    ("Resellers", "rank hundreds of listings by expected profit in seconds, never overpay.", GREEN),
    ("Buyers", "an objective fair-value second opinion + a scam radar before you message a seller.", BLUE),
    ("Researchers", "a reproducible, versioned pipeline for second-hand price dynamics.", PURPLE)]):
    x = Inches(0.7 + i*4.05)
    g = glass(s, x, Inches(1.9), Inches(3.8), Inches(2.7), a=9)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.25), Inches(2.15), Inches(0.5), Inches(0.12))
    _alpha(strip, c, 100); _no_line(strip); _noshadow(strip)
    text(s, x+Inches(0.25), Inches(2.4), Inches(3.3), Inches(2), [(t, 18, True, WHITE, HEAD), ("", 4, False, SUB, BODY), (b, 13.5, False, SUB, BODY)], sp=6)
g = glass(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.8), a=7)
bullets(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.4), [
    ("vs manual", "objective (model, not gut) · exhaustive (every listing) · fast (~1,670 listings/sec) · uncertainty-aware."),
    ("vs naive baseline", "Ridge floor R² 0.58 → boosting 0.64; calibrated intervals turn a guess into a confidence."),
], size=14.5, gap=10)

# 11 · STRUGGLES (Rubric 5)
s = slide(); kicker(s, "05 · Struggles & Problem-Solving /15", BLUE); title(s, "Six real problems, six fixes")
items = [
    ("Quantile intervals under-covered (60%)", "Conformalized Quantile Regression → ~80%."),
    ("eBay/Backmarket returned HTTP 403", "curl_cffi TLS (JA3) impersonation + cookie-warming."),
    ("eBay's sold search is capped", "diagnosed; scaled via narrow model×storage queries."),
    ("Merging US + ES data dropped R² 0.74→0.61", "added market & source features; reweighted to the ES market."),
    ("SHAP crashed on XGBoost 2.x", "used XGBoost's native tree-SHAP (pred_contribs)."),
    ("Accessories/parts polluted the data", "title junk filter + damaged-device risk flag."),
]
g = glass(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.8), a=7)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(4.4)); tf = tb.text_frame; tf.word_wrap = True
for i, (p_, f_) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(13)
    r = p.add_run(); r.text = "● "; r.font.color.rgb = AMBER; r.font.size = Pt(15)
    r1 = p.add_run(); r1.text = p_ + "   →   "; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = BODY
    r2 = p.add_run(); r2.text = f_; r2.font.size = Pt(15); r2.font.color.rgb = GREEN; r2.font.name = BODY

# 12 · COURSE CONCEPTS + feature importance
s = slide(); kicker(s, "What we used from the course"); title(s, "Concepts → code")
mods = [("S1 OOP", "abstract classes · __slots__ · exceptions"), ("S2 Performance", "parquet · vectorized · timeit (5.9× CSV)"),
        ("S3 Visualization", "7 figure types"), ("S4 Parallel & Exceptions", "retry/backoff · custom errors"),
        ("S5 Web Scraping", "JSON API · sold HTML · TLS impersonation"), ("S6 Time Series", "out-of-time split · TimeSeriesSplit"),
        ("S7 Supervised ML", "Ridge·XGBoost·LightGBM · CQR · SHAP")]
y = Inches(1.95)
for i, (t, b) in enumerate(mods):
    col = i % 2; x = Inches(0.7 + col*3.55)
    yy = Emu(int(y) + (i // 2) * int(Inches(0.92)))
    g = glass(s, x, yy, Inches(3.4), Inches(0.82), a=9)
    text(s, x+Inches(0.2), yy+Inches(0.08), Inches(3.0), Inches(0.7), [(t, 13.5, True, WHITE, HEAD), (b, 10.5, False, SUB, BODY)], sp=1)
pic(s, "slide_feature_importance.png", Inches(7.7), Inches(1.95), w=Inches(5.0))
text(s, Inches(7.7), Inches(6.5), Inches(5.0), Inches(0.5), [("36 tests · 4 notebooks · modular package", 12.5, True, SUB, MONO)], align=PP_ALIGN.CENTER, sp=0)

# 13 · NEXT STEPS + close (Rubric 6)
s = slide(); glow(s, int(SW*0.5), int(SH*0.4), Inches(8), "0A84FF", 10)
kicker(s, "06 · Next Steps /10", BLUE); title(s, "Where this goes next")
g = glass(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.4), a=8)
bullets(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(3.0), [
    ("Counterfeit detection", "a CNN on listing photos to flag probable fakes (S7 deep learning)."),
    ("Time-to-sell", "predict liquidity (how fast a deal sells) as a 4th score dimension."),
    ("True panel data", "scheduled daily scrape to detect actual sales, price drops, time-on-market."),
    ("Productionize", "FastAPI serve mode, CI on the test suite, GA hyperparameter search vs grid."),
], size=15.5, gap=13)
text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.8), [("Live data → fair value → a confident buy/skip call.", 22, True, WHITE, HEAD)], sp=2)
text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6),
     [("Academic project · the investment score is not financial advice · do your own due diligence.", 12, False, SUB, BODY)], sp=0)

out = ROOT / "iPhone_Deal_Finder.pptx"
prs.save(out)
print("saved", out, "| slides:", len(prs.slides._sldIdLst))
